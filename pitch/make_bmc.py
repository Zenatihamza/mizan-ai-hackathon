"""Generate a single-slide Business Model Canvas for JusticIA (16:9 PPTX)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# JusticIA warm palette
ESPRESSO = RGBColor(0x2C, 0x1C, 0x14)
MAROON = RGBColor(0x7A, 0x2E, 0x2A)
COPPER = RGBColor(0xC8, 0x79, 0x4F)
GOLD = RGBColor(0xA8, 0x76, 0x3A)
SAND = RGBColor(0xF5, 0xEC, 0xE2)
BORDER = RGBColor(0xD8, 0xC4, 0xB4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5C, 0x42, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.28), Inches(12.5), Inches(0.95))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Business Model Canvas"
r.font.size = Pt(30); r.font.bold = True; r.font.name = "Calibri"; r.font.color.rgb = ESPRESSO
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "JusticIA  ·  l'accès au droit algérien, en français et en arabe  ·  Hackathon Nous — Challenge 1"
r2.font.size = Pt(13); r2.font.name = "Calibri"; r2.font.color.rgb = MAROON


def block(x, y, w, h, title, bullets, fill, title_color, body_color, cols=1):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = BORDER; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.07); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    rt = p.add_run(); rt.text = title.upper()
    rt.font.size = Pt(11.5); rt.font.bold = True; rt.font.name = "Calibri"
    rt.font.color.rgb = title_color
    p.space_after = Pt(4)
    half = (len(bullets) + 1) // 2 if cols == 2 else len(bullets)
    for i, b in enumerate(bullets):
        bp = tf.add_paragraph()
        rb = bp.add_run(); rb.text = "•  " + b
        rb.font.size = Pt(10.5); rb.font.name = "Calibri"; rb.font.color.rgb = body_color
        bp.space_after = Pt(2)
    return shp


# Geometry
cw, gap = 2.4106, 0.12
x = [0.4 + i * (cw + gap) for i in range(5)]
top_y = 1.4
top_h = 4.16
half_h = (top_h - 0.12) / 2
bot_y = top_y + top_h + 0.12
bot_h = 7.2 - bot_y

block(x[0], top_y, cw, top_h, "Partenaires clés",
      ["Ministère de la Justice, Dzair Digital", "Avocats & aide juridictionnelle",
       "Associations (APOCE, consommateurs)", "Universités & facultés de droit",
       "Maisons de justice, mairies"], SAND, MAROON, ESPRESSO)

block(x[1], top_y, cw, half_h, "Activités clés",
      ["Développement de la plateforme", "Curation du corpus juridique",
       "Entraînement du modèle ML"], SAND, MAROON, ESPRESSO)
block(x[1], top_y + half_h + 0.12, cw, half_h, "Ressources clés",
      ["Corpus de lois algériennes", "Modèle ML entraîné",
       "Équipe tech + juridique"], SAND, MAROON, ESPRESSO)

block(x[2], top_y, cw, top_h, "Proposition de valeur",
      ["Comprendre ses droits en FR & arabe", "Réponses ancrées dans la loi, sans invention",
       "Scanner de clauses abusives", "Mode urgence + démarches par wilaya",
       "Gratuit, accessible, sans avocat"], COPPER, WHITE, WHITE)

block(x[3], top_y, cw, half_h, "Relations clients",
      ["Assistant en libre-service", "Gamification (simulateur)",
       "Escalade vers de vrais avocats"], SAND, MAROON, ESPRESSO)
block(x[3], top_y + half_h + 0.12, cw, half_h, "Canaux",
      ["Application web puis mobile", "Universités, réseaux sociaux",
       "Guichets publics partenaires"], SAND, MAROON, ESPRESSO)

block(x[4], top_y, cw, top_h, "Segments de clientèle",
      ["Citoyens, populations vulnérables", "Ruraux, faibles revenus",
       "Peu alphabétisés", "Étudiants & universités", "Associations & ONG"],
      SAND, MAROON, ESPRESSO)

cost_w = 3 * cw + 2 * gap
block(0.4, bot_y, cost_w, bot_h, "Structure de coûts",
      ["Développement & hébergement", "Juristes pour la curation",
       "Marketing & partenariats", "Coût léger : CPU, open-source"],
      SAND, MAROON, ESPRESSO, cols=2)

rev_w = 2 * cw + gap
block(x[3], bot_y, rev_w, bot_h, "Sources de revenus",
      ["Freemium, gratuit pour les citoyens", "B2B : cabinets & associations",
       "Subventions & partenariats publics", "Premium : mise en relation, rapports PDF"],
      SAND, GOLD, ESPRESSO, cols=2)

prs.save("C:/Users/Hello/Desktop/mizan-ai/pitch/JusticIA_BMC.pptx")
print("saved JusticIA_BMC.pptx")
